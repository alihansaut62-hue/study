import sys
from pptx import Presentation

def extract_text(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = []
        for i, slide in enumerate(prs.slides):
            text.append(f"--- Слайд {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        with open("pptx_text.txt", "w", encoding="utf-8") as f:
            f.write("\n".join(text))
    except Exception as e:
        print("Error:", e)

if __name__ == "__main__":
    extract_text(sys.argv[1])
